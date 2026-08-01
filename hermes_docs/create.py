"""hermes_plugins.doc_create — Create .docx, .xlsx, .pptx, .pdf, .odt, .ods, .odp documents from structured JSON.

Uses a dedicated venv at ~/.hermes/venvs/docs with python-docx, openpyxl,
python-pptx, weasyprint, odfpy, Pillow, and lxml.
"""

import json
import logging
import os
import platform
import shutil
import subprocess
import tempfile
from pathlib import Path

logger = logging.getLogger(__name__)

from hermes_docs import get_docs_venv_python
VENV_PYTHON = get_docs_venv_python()


def _check_weasyprint_deps() -> str | None:
    """Check if WeasyPrint's system dependencies are available.

    WeasyPrint needs cairo, pango, and gdk-pixbuf C libraries. On Linux
    these are typically installed via the package manager. On macOS via brew.

    Returns an error string if deps are missing, None if OK.
    """
    try:
        # Try importing weasyprint in the docs venv — if it imports, the
        # C libs are present and we're good.
        result = subprocess.run(
            [VENV_PYTHON, "-c", "import weasyprint; print('ok')"],
            capture_output=True, text=True, timeout=10,
        )
        if result.returncode == 0:
            return None
        # Import failed — likely missing system C libs
        system = platform.system()
        if system == "Darwin":
            hint = "brew install cairo pango gdk-pixbuf"
        elif system == "Linux":
            if shutil.which("pacman"):
                hint = "sudo pacman -S cairo pango gdk-pixbuf2"
            elif shutil.which("dnf"):
                hint = "sudo dnf install cairo pango gdk-pixbuf2"
            elif shutil.which("apt"):
                hint = "sudo apt install libcairo2 libpango-1.0-0 libgdk-pixbuf-2.0-0"
            else:
                hint = "install cairo, pango, and gdk-pixbuf via your package manager"
        else:
            hint = "install cairo, pango, and gdk-pixbuf (see https://doc.courtbouillon.org/weasyprint/stable/first_steps.html)"
        return f"WeasyPrint system dependencies missing. Install: `{hint}`"
    except Exception as exc:
        return f"Could not verify WeasyPrint dependencies: {exc}"


def _run_in_venv(script: str) -> dict:
    """Execute a Python script inside the docs venv and return JSON result."""
    result = subprocess.run(
        [VENV_PYTHON, "-c", script],
        capture_output=True, text=True, timeout=120
    )
    if result.returncode != 0:
        return {"success": False, "error": result.stderr.strip()}
    try:
        return json.loads(result.stdout.strip())
    except json.JSONDecodeError:
        return {"success": False, "error": f"Invalid JSON from subprocess: {result.stdout[:500]}"}


# ── builder script generators ──

def _build_docx_script(content: dict, output_path: str, template: str = None, options: dict = None) -> str:
    content_json = json.dumps(content)
    options_json = json.dumps(options or {})
    template_str = template or ""
    script = f'''
import json, os
from docx import Document
from docx.shared import Inches

content = json.loads({repr(content_json)})
options = json.loads({repr(options_json)})
template_path = {repr(template_str)}
output_path = {repr(output_path)}

try:
    if template_path and os.path.exists(template_path):
        doc = Document(template_path)
    else:
        doc = Document()

    if options.get("title"):
        doc.add_heading(options["title"], level=0)

    for section in content.get("sections", []):
        stype = section.get("type", "paragraph")
        if stype == "heading":
            doc.add_heading(section.get("text", ""), level=section.get("level", 1))
        elif stype == "paragraph":
            p = doc.add_paragraph(section.get("text", ""))
            if section.get("style"):
                p.style = section["style"]
        elif stype == "table":
            rows = section.get("rows", [])
            headers = section.get("headers", [])
            total_rows = len(rows) + (1 if headers else 0)
            num_cols = max(1, len(headers) or (len(rows[0]) if rows else 1))
            table = doc.add_table(rows=max(1, total_rows), cols=num_cols)
            table.style = "Table Grid"
            row_idx = 0
            if headers:
                for col_idx, h in enumerate(headers):
                    table.rows[row_idx].cells[col_idx].text = str(h)
                row_idx += 1
            for row_data in rows:
                for col_idx, cell in enumerate(row_data):
                    if col_idx < num_cols:
                        table.rows[row_idx].cells[col_idx].text = str(cell)
                row_idx += 1
        elif stype == "image":
            doc.add_picture(section["path"], width=Inches(section.get("width", 5)))
        elif stype == "page_break":
            doc.add_page_break()

    doc.save(output_path)
    print(json.dumps({{"success": True, "path": output_path, "format": "docx"}}))
except Exception as e:
    print(json.dumps({{"success": False, "error": str(e)}}))
'''
    return script


def _build_xlsx_script(content: dict, output_path: str, template: str = None, options: dict = None) -> str:
    content_json = json.dumps(content)
    options_json = json.dumps(options or {})
    script = f'''
import json, os
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill
from openpyxl.utils import get_column_letter

content = json.loads({repr(content_json)})
options = json.loads({repr(options_json)})
output_path = {repr(output_path)}

try:
    wb = Workbook()
    for sheet in content.get("sheets", []):
        ws = wb.active if sheet == content["sheets"][0] else wb.create_sheet(title=sheet.get("name", "Sheet"))
        if "name" in sheet and sheet != content["sheets"][0]:
            ws.title = sheet["name"]

        data = sheet.get("data", [])
        formulas = sheet.get("formulas", {{}})
        formats = sheet.get("formats", {{}})

        for row_idx, row in enumerate(data, 1):
            for col_idx, val in enumerate(row, 1):
                cell_ref = f"{{get_column_letter(col_idx)}}{{row_idx}}"
                cell = ws.cell(row=row_idx, column=col_idx, value=val)
                if formats.get(cell_ref) == "header" or formats.get(cell_ref) == "bold":
                    cell.font = Font(bold=True)
                if cell_ref in formulas:
                    cell.value = formulas[cell_ref]

        if options.get("auto_size", True):
            for col in ws.columns:
                max_length = 0
                column = col[0].column_letter
                for cell in col:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                ws.column_dimensions[column].width = adjusted_width

    wb.save(output_path)
    print(json.dumps({{"success": True, "path": output_path, "format": "xlsx"}}))
except Exception as e:
    print(json.dumps({{"success": False, "error": str(e)}}))
'''
    return script


def _build_pptx_script(content: dict, output_path: str, template: str = None, options: dict = None) -> str:
    content_json = json.dumps(content)
    options_json = json.dumps(options or {})
    script = f'''
import json, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

content = json.loads({repr(content_json)})
options = json.loads({repr(options_json)})
output_path = {repr(output_path)}
template_path = {repr(template)}

try:
    prs = Presentation(template_path) if template_path and os.path.exists(template_path) else Presentation()

    for slide in content.get("slides", []):
        layout_name = slide.get("layout", "title_and_content")
        layout_idx = 0
        for i, l in enumerate(prs.slide_layouts):
            if l.name.lower().replace(" ", "_") == layout_name.lower().replace(" ", "_"):
                layout_idx = i
                break
        s = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        if slide.get("title"):
            s.shapes.title.text = slide["title"]
        if slide.get("subtitle"):
            for shape in s.placeholders:
                if shape.placeholder_format.type == 2:
                    shape.text = slide["subtitle"]
                    break
        if slide.get("bullets"):
            body = s.shapes.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for i, bullet in enumerate(slide["bullets"]):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
        if slide.get("table"):
            tbl_data = slide["table"]
            rows = len(tbl_data.get("data", [])) + (1 if tbl_data.get("headers") else 0)
            cols = max(len(tbl_data.get("headers", [])), len(tbl_data.get("data", [[]])[0]) if tbl_data.get("data") else 1)
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(2)
            table = s.shapes.add_table(rows, cols, left, top, width, height).table
            if tbl_data.get("headers"):
                for i, h in enumerate(tbl_data["headers"]):
                    table.cell(0, i).text = str(h)
            for r_idx, row in enumerate(tbl_data.get("data", [])):
                for c_idx, cell in enumerate(row):
                    table.cell(r_idx + (1 if tbl_data.get("headers") else 0), c_idx).text = str(cell)

    prs.save(output_path)
    print(json.dumps({{"success": True, "path": output_path, "format": "pptx"}}))
except Exception as e:
    print(json.dumps({{"success": False, "error": str(e)}}))
'''
    return script


def _build_pdf_script(content: dict, output_path: str, template: str = None, options: dict = None) -> str:
    # Auto-wrap string content into structured format for convenience
    if isinstance(content, str):
        content = {
            "pages": [
                {
                    "elements": [
                        {"type": "heading", "text": content}
                    ]
                }
            ]
        }
    content_json = json.dumps(content)
    options_json = json.dumps(options or {})
    script = f'''
import json, os
from weasyprint import HTML, CSS
from weasyprint.text.fonts import FontConfiguration

content = json.loads({repr(content_json)})
options = json.loads({repr(options_json)})
output_path = {repr(output_path)}

try:
    page_size = options.get("page_size", "letter")
    title = options.get("title", "")
    header_text = options.get("header", "")
    footer_text = options.get("footer", "")
    page_numbers = options.get("page_numbers", True)

    css_size = "letter" if page_size.lower() == "letter" else "A4"

    html_parts = ['<html><head><meta charset="utf-8"><title>' + (title or "Document") + '</title></head><body>']

    for page in content.get("pages", []):
        html_parts.append('<div class="page">')
        for el in page.get("elements", []):
            etype = el.get("type", "text")
            if etype == "title":
                html_parts.append(f"<h1>{{el.get('text', '')}}</h1>")
            elif etype == "heading":
                html_parts.append(f"<h2>{{el.get('text', '')}}</h2>")
            elif etype == "text":
                html_parts.append(f"<p>{{el.get('text', '')}}</p>")
            elif etype == "table":
                rows = el.get("rows", [])
                html_parts.append('<table class="data-table">')
                for row in rows:
                    html_parts.append("<tr>")
                    for cell in row:
                        html_parts.append(f"<td>{{cell}}</td>")
                    html_parts.append("</tr>")
                html_parts.append("</table>")
            elif etype == "spacer":
                h = el.get("height", 24)
                html_parts.append(f'<div style="height:{{h}}px"></div>')
        html_parts.append('</div>')

    html_parts.append('</body></html>')
    html_str = "\\n".join(html_parts)

    css_str = """
        @page {{ size: """ + css_size + """; margin: 1in; }}
        body {{ font-family: Helvetica, Arial, sans-serif; font-size: 11pt; line-height: 1.4; }}
        h1 {{ font-size: 18pt; margin-bottom: 12pt; }}
        h2 {{ font-size: 14pt; margin-top: 14pt; margin-bottom: 8pt; }}
        table {{ width: 100%; border-collapse: collapse; margin: 10pt 0; }}
        td, th {{ border: 1px solid #ccc; padding: 6px; text-align: left; }}
        .page {{ page-break-after: always; }}
        .page:last-child {{ page-break-after: auto; }}
    """

    if header_text:
        css_str += """
            @page {{ @top-center {{ content: '" + header_text + "'; font-size: 9pt; color: #666; }} }}
        """
    if footer_text or page_numbers:
        footer_content = footer_text or ""
        if page_numbers:
            footer_content += " Page " + ("" if not footer_text else " ") + "counter(page)"
        css_str += """
            @page {{ @bottom-center {{ content: '" + footer_content + "'; font-size: 9pt; color: #666; }} }}
        """

    font_config = FontConfiguration()
    HTML(string=html_str).write_pdf(output_path, stylesheets=[CSS(string=css_str, font_config=font_config)])
    print(json.dumps({{"success": True, "path": output_path, "format": "pdf"}}))
except Exception as e:
    print(json.dumps({{"success": False, "error": str(e)}}))
'''
    return script


def _build_odt_script(content: dict, output_path: str, template: str = None, options: dict = None) -> str:
    content_json = json.dumps(content)
    options_json = json.dumps(options or {})
    script = f'''
import json
from odf.opendocument import OpenDocumentText
from odf.style import Style, TextProperties, ParagraphProperties
from odf.text import P, H, TableOfContent
from odf.table import Table, TableRow, TableCell

content = json.loads({repr(content_json)})
options = json.loads({repr(options_json)})
output_path = {repr(output_path)}

try:
    doc = OpenDocumentText()
    if options.get("title"):
        h = H(outlinelevel=1, text=options["title"])
        doc.text.addElement(h)

    for section in content.get("sections", []):
        stype = section.get("type", "paragraph")
        if stype == "heading":
            h = H(outlinelevel=section.get("level", 1), text=section.get("text", ""))
            doc.text.addElement(h)
        elif stype == "paragraph":
            p = P(text=section.get("text", ""))
            doc.text.addElement(p)
        elif stype == "table":
            rows = section.get("rows", [])
            headers = section.get("headers", [])
            total_rows = len(rows) + (1 if headers else 0)
            num_cols = max(1, len(headers) or (len(rows[0]) if rows else 1))
            table = Table()
            for r in range(total_rows):
                tr = TableRow()
                for c in range(num_cols):
                    tc = TableCell()
                    p = P(text="")
                    tc.addElement(p)
                    tr.addElement(tc)
                table.addElement(tr)
            doc.text.addElement(table)

    doc.save(output_path)
    print(json.dumps({{"success": True, "path": output_path, "format": "odt"}}))
except Exception as e:
    print(json.dumps({{"success": False, "error": str(e)}}))
'''
    return script


def _build_ods_script(content: dict, output_path: str, template: str = None, options: dict = None) -> str:
    content_json = json.dumps(content)
    script = f'''
import json
from odf.opendocument import OpenDocumentSpreadsheet
from odf.table import Table, TableRow, TableCell
from odf.text import P

content = json.loads({repr(content_json)})
output_path = {repr(output_path)}

try:
    doc = OpenDocumentSpreadsheet()
    for sheet in content.get("sheets", []):
        table = Table(name=sheet.get("name", "Sheet1"))
        data = sheet.get("data", [])
        for row in data:
            tr = TableRow()
            for cell in row:
                tc = TableCell()
                p = P(text=str(cell))
                tc.addElement(p)
                tr.addElement(tc)
            table.addElement(tr)
        doc.spreadsheet.addElement(table)
    doc.save(output_path)
    print(json.dumps({{"success": True, "path": output_path, "format": "ods"}}))
except Exception as e:
    print(json.dumps({{"success": False, "error": str(e)}}))
'''
    return script


def _build_odp_script(content: dict, output_path: str, template: str = None, options: dict = None) -> str:
    content_json = json.dumps(content)
    script = f'''
import json
from odf.opendocument import OpenDocumentPresentation
from odf.style import Style, GraphicProperties, PresentationPageLayout
from odf.draw import Page, Frame, TextBox
from odf.text import P

content = json.loads({repr(content_json)})
output_path = {repr(output_path)}

try:
    doc = OpenDocumentPresentation()
    for slide in content.get("slides", []):
        page = Page(masterpagename="Default")
        doc.presentation.addElement(page)
        if slide.get("title"):
            frame = Frame(x="1cm", y="1cm", width="18cm", height="2cm")
            textbox = TextBox()
            p = P(text=slide["title"])
            textbox.addElement(p)
            frame.addElement(textbox)
            page.addElement(frame)
        if slide.get("subtitle"):
            frame = Frame(x="1cm", y="3.5cm", width="18cm", height="1cm")
            textbox = TextBox()
            p = P(text=slide["subtitle"])
            textbox.addElement(p)
            frame.addElement(textbox)
            page.addElement(frame)
        if slide.get("bullets"):
            frame = Frame(x="1cm", y="5cm", width="18cm", height="10cm")
            textbox = TextBox()
            for bullet in slide["bullets"]:
                p = P(text="• " + bullet)
                textbox.addElement(p)
            frame.addElement(textbox)
            page.addElement(frame)
    doc.save(output_path)
    print(json.dumps({{"success": True, "path": output_path, "format": "odp"}}))
except Exception as e:
    print(json.dumps({{"success": False, "error": str(e)}}))
'''
    return script


# ── content normalization ──

def _is_html(text: str) -> bool:
    """Heuristic: does the string look like HTML?"""
    if not isinstance(text, str):
        return False
    stripped = text.strip()
    return stripped.startswith("<") and (">" in stripped or stripped.lower().startswith("<!doctype"))


def _normalize_content(format: str, content, options: dict = None) -> dict:
    """Convert string content into the structured dict expected by each builder.

    If content is already a dict, return it as-is.
    If content is a plain string:
      - For PDF: if it looks like HTML, wrap as {'html': '<raw html>'};
        otherwise wrap as a single text element page.
      - For DOCX/ODT: wrap as a single paragraph section.
      - For XLSX/ODS: wrap as a single-cell single-sheet.
      - For PPTX/ODP: wrap as a single title slide.
    """
    if isinstance(content, dict):
        return content

    if not isinstance(content, str):
        content = str(content)

    fmt = format.lower()

    if fmt == "pdf":
        if _is_html(content):
            return {"html": content}
        return {
            "pages": [
                {
                    "elements": [
                        {"type": "text", "text": content}
                    ]
                }
            ]
        }

    if fmt in ("docx", "odt"):
        return {
            "sections": [
                {"type": "paragraph", "text": content}
            ]
        }

    if fmt in ("xlsx", "ods"):
        return {
            "sheets": [
                {
                    "name": "Sheet1",
                    "data": [[content]]
                }
            ]
        }

    if fmt in ("pptx", "odp"):
        return {
            "slides": [
                {"title": content, "layout": "title_slide"}
            ]
        }

    return content


# ── raw HTML PDF builder ──

def _build_pdf_from_html_script(content: dict, output_path: str, template: str = None, options: dict = None) -> str:
    """Build a PDF directly from raw HTML string."""
    html_str = content.get("html", "")
    options_json = json.dumps(options or {})
    script = f'''
import json, os
from weasyprint import HTML, CSS
from weasyprint.text.fonts import FontConfiguration

output_path = {repr(output_path)}
options = json.loads({repr(options_json)})
html_str = {repr(html_str)}

try:
    page_size = options.get("page_size", "letter")
    css_size = "letter" if page_size.lower() == "letter" else "A4"
    title = options.get("title", "")
    header_text = options.get("header", "")
    footer_text = options.get("footer", "")
    page_numbers = options.get("page_numbers", True)

    if "<!doctype" not in html_str.lower() and "<html" not in html_str.lower():
        _t = title or "Document"
        html_str = (
            "\\n<!DOCTYPE html>\\n<html><head><meta charset=\\\"utf-8\\\"><title>" + _t +
            "</title></head><body>\\n" + html_str + "</body></html>"
        )

    css_lines = []
    css_lines.append("@page {{ size: " + css_size + "; margin: 1in; }}")
    css_lines.append("body {{ font-family: Helvetica, Arial, sans-serif; font-size: 11pt; line-height: 1.4; }}")
    css_lines.append("h1 {{ font-size: 18pt; margin-bottom: 12pt; }}")
    css_lines.append("h2 {{ font-size: 14pt; margin-top: 14pt; margin-bottom: 8pt; }}")
    css_lines.append("table {{ width: 100%; border-collapse: collapse; margin: 10pt 0; }}")
    css_lines.append("td, th {{ border: 1px solid #ccc; padding: 6px; text-align: left; }}")

    if header_text:
        css_lines.append("@page {{ @top-center {{ content: '" + header_text.replace("'", "\\\\'") + "'; font-size: 9pt; color: #666; }} }}")
    if footer_text or page_numbers:
        footer_content = footer_text or ""
        if page_numbers:
            footer_content += " Page " + ("" if not footer_text else " ") + "counter(page)"
        css_lines.append("@page {{ @bottom-center {{ content: '" + footer_content.replace("'", "\\\\'") + "'; font-size: 9pt; color: #666; }} }}")

    css_str = "\\n".join(css_lines)
    font_config = FontConfiguration()
    HTML(string=html_str).write_pdf(output_path, stylesheets=[CSS(string=css_str, font_config=font_config)])
    print(json.dumps({{"success": True, "path": output_path, "format": "pdf"}}))
except Exception as e:
    print(json.dumps({{"success": False, "error": str(e)}}))
'''
    return script


# ── core create_document ──

def _get_user_dir() -> Path:
    """Resolve output directory for documents."""
    from hermes_docs import get_output_dir
    return get_output_dir()


def _is_pcloud_path(path: Path) -> bool:
    """Check if path is under the configured output directory."""
    try:
        resolved = path.expanduser().resolve()
        user_dir = _get_user_dir()
        return str(resolved).startswith(str(user_dir.parent))
    except Exception:
        return False


def create_document(format: str, content, output_path: str, template: str = None, options: dict = None) -> dict:
    """Create a document in the specified format.

    Args:
        format: "docx", "xlsx", "pptx", "pdf", "odt", "ods", or "odp"
        content: Structured content dict, or a plain string (auto-normalized)
        output_path: Absolute path for the output file
        template: Optional template path
        options: Optional format-specific options
    """
    format = format.lower()
    out = Path(output_path).expanduser().resolve()
    out.parent.mkdir(parents=True, exist_ok=True)

    if not _is_pcloud_path(out):
        logger.warning(
            "doc_create: output file %s is outside the configured output directory.",
            out,
        )

    builders = {
        "docx": _build_docx_script,
        "xlsx": _build_xlsx_script,
        "pptx": _build_pptx_script,
        "pdf": _build_pdf_script,
        "odt": _build_odt_script,
        "ods": _build_ods_script,
        "odp": _build_odp_script,
    }

    if format not in builders:
        return {"success": False, "error": f"Unsupported format: {format}"}

    if not shutil.which(VENV_PYTHON):
        return {"success": False, "error": f"docs venv not found at {VENV_PYTHON}. Create it with: python -m venv ~/.hermes/venvs/docs && ~/.hermes/venvs/docs/bin/pip install python-docx openpyxl python-pptx weasyprint odfpy Pillow lxml PyMuPDF"}

    # WeasyPrint needs cairo/pango system libs — check when PDF is requested
    if format == "pdf":
        wp_error = _check_weasyprint_deps()
        if wp_error:
            return {"success": False, "error": wp_error}

    # Normalize string content into structured dict
    normalized = _normalize_content(format, content, options=options)

    # PDF: if normalized content has raw HTML key, use the HTML-specific builder
    if format == "pdf" and "html" in normalized:
        script = _build_pdf_from_html_script(normalized, str(out), template=template, options=options or {})
    else:
        script = builders[format](normalized, str(out), template=template, options=options or {})

    result = _run_in_venv(script)
    if result.get("success") and "path" in result:
        result["message"] = f"MEDIA:{result['path']}\n\nDocument created ({format}): {result['path']}"
    return result


# ── PDF editing ──

def _edit_pdf_text(path: str, page: int, instruction: str, output_path: str = None) -> dict:
    """Edit PDF text via pymupdf search-and-replace (no LLM, provider-agnostic)."""
    try:
        import fitz  # pymupdf
    except ImportError:
        return {"success": False, "error": "pymupdf not installed."}

    doc = fitz.open(path)
    if page < 1 or page > len(doc):
        return {"success": False, "error": f"Page {page} out of range (1-{len(doc)})."}

    page_obj = doc[page - 1]
    original_text = page_obj.get_text()

    # Try to parse "change X to Y" / "replace X with Y" / "fix X to Y" patterns
    import re
    patterns = [
        r"(?:change|replace|fix|update|correct)\s+['\"]?(.+?)['\"]?\s+(?:to|with|->)\s+['\"]?(.+?)['\"]?$",
        r"['\"]?(.+?)['\"]?\s*->\s*['\"]?(.+?)['\"]?$",
    ]

    search_text = None
    replace_text = None
    for pattern in patterns:
        m = re.search(pattern, instruction.strip(), re.IGNORECASE)
        if m:
            search_text = m.group(1).strip().strip("'\"")
            replace_text = m.group(2).strip().strip("'\"")
            break

    if not search_text:
        return {
            "success": False,
            "error": (
                "Could not parse instruction. Use format: "
                "'change <search> to <replace>' or 'replace <search> with <replace>'"
            ),
            "page_text": original_text[:2000],
        }

    # pymupdf text replacement: find rects, redact, insert new text
    found = False
    text_instances = page_obj.search_for(search_text)

    if not text_instances:
        # Try case-insensitive fallback by scanning all text blocks
        blocks = page_obj.get_text("blocks")
        for block in blocks:
            if search_text.lower() in block[4].lower():
                # Found in a block, try to locate it more precisely
                text_instances = page_obj.search_for(block[4].strip()[:len(search_text) + 20])
                break

    if not text_instances:
        return {
            "success": False,
            "error": f"Text '{search_text}' not found on page {page}.",
            "page_text": original_text[:2000],
        }

    for rect in text_instances:
        # Redact the area and insert replacement text
        page_obj.add_redact_annot(rect, text=replace_text, fontsize=11)
        found = True

    if found:
        page_obj.apply_redactions()

    if not output_path:
        output_path = str(Path(path).parent / f"{Path(path).stem}_edited.pdf")

    doc.save(output_path)
    doc.close()

    return {
        "success": True,
        "path": output_path,
        "format": "pdf",
        "page": page,
        "instruction": instruction,
        "search": search_text,
        "replace": replace_text,
        "message": f"PDF edited (page {page}). Output: {output_path}",
    }


# Kept for backward compat — calls the text-based editor
_edit_pdf_via_nano_pdf = _edit_pdf_text


# ── Hermes handlers ──

def _handle_doc_create(params: dict, **kwargs) -> str:
    """Handler for doc_create_document tool."""
    # Accept either "output_path" or "path" for backward compatibility
    output_path = params.get("output_path") or params.get("path")
    
    # Default output directory with timestamped filename
    import time
    cache_dir = _get_user_dir()
    cache_dir.mkdir(parents=True, exist_ok=True)
    
    if output_path and "/tmp/" in output_path:
        # Redirect /tmp paths to cache/documents for media delivery compatibility
        ts = int(time.time())
        fmt = params["format"]
        p = Path(output_path)
        if p.name and p.name != ".":
            filename = p.stem + f"_{ts}.{fmt}"
        else:
            filename = f"doc_{ts}.{fmt}"
        output_path = str(cache_dir / filename)
        logger.warning("doc_create: redirecting /tmp path to %s for media delivery compatibility", output_path)
    if not output_path:
        ts = int(time.time())
        fmt = params["format"]
        output_path = str(cache_dir / f"doc_{ts}.{fmt}")
    result = create_document(
        format=params["format"],
        content=params["content"],
        output_path=output_path,
        template=params.get("template"),
        options=params.get("options"),
    )
    if isinstance(result, dict):
        return json.dumps(result)
    return result


def _handle_doc_edit(params: dict, **kwargs) -> str:
    """Handler for doc_edit_pdf tool."""
    result = _edit_pdf_via_nano_pdf(
        path=params["path"],
        page=params.get("page", 1),
        instruction=params["instruction"],
        output_path=params.get("output_path"),
    )
    if isinstance(result, dict):
        return json.dumps(result)
    return result


# ── plugin registration ──
def register(registry):
    registry.register_tool(
        name="doc_create_document",
        toolset="doc-create",
        schema={
            "type": "object",
            "properties": {
                "format": {
                    "type": "string",
                    "enum": ["docx", "xlsx", "pptx", "pdf", "odt", "ods", "odp"],
                    "description": "Output document format.",
                },
                "content": {
                    "anyOf": [
                        {"type": "string", "description": "Plain text or HTML content (auto-detected and normalized)."},
                        {"type": "object", "description": "Structured content dict. Format-specific schema (see references/)."},
                    ],
                    "description": "Document content. Accepts plain text, HTML, or structured dict.",
                },
                "output_path": {
                    "type": "string",
                    "description": "Absolute path for the output file. Also accepts 'path' as an alias.",
                },
                "path": {
                    "type": "string",
                    "description": "Alias for output_path. Absolute path for the output file.",
                },
                "template": {
                    "type": "string",
                    "description": "Optional template path (.dotx, .odt, .xltx, .potx, etc.)",
                },
                "options": {
                    "type": "object",
                    "description": "Optional format-specific options (title, page_size, auto_size, etc.)",
                },
            },
            "required": ["format", "content"],
        },
        handler=_handle_doc_create,
        description="Create .docx, .xlsx, .pptx, .pdf, .odt, .ods, or .odp documents. For PDFs: pass structured JSON for document-style output, or pass {\"html\": \"<html>...</html>\"} with full CSS (Grid, Flexbox, gradients, themes) for visual/print-ready PDFs via weasyprint. Supports page_size, header, footer, page_numbers options.",
    )

    registry.register_tool(
        name="doc_edit_pdf",
        toolset="doc-create",
        schema={
            "type": "object",
            "properties": {
                "path": {
                    "type": "string",
                    "description": "Absolute path to the PDF file.",
                },
                "page": {
                    "type": "integer",
                    "description": "Page number to edit (1-indexed).",
                    "default": 1,
                },
                "instruction": {
                    "type": "string",
                    "description": "Natural language instruction for the edit.",
                },
                "output_path": {
                    "type": "string",
                    "description": "Optional output path. Defaults to <original>_edited.pdf",
                },
            },
            "required": ["path", "instruction"],
        },
        handler=_handle_doc_edit,
        description="Edit a PDF page using AI-powered visual editing (nano-pdf + Gemini).",
    )
